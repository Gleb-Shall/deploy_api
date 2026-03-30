"""
Конвертация файлов:
  PDF  → список PNG (по одному на страницу) через pymupdf
  DOCX → текст (.txt)            через python-docx
  PPTX → текст (.txt)            через python-pptx
  XLSX → текст (.txt)            через openpyxl
  XLS  → текст (.txt)            через xlrd
"""
import io


class ProcessingError(Exception):
    """Ошибка конвертации файла с сообщением для пользователя."""


def pdf_to_pages(content: bytes, max_pages: int = 50) -> list[tuple[bytes, str]]:
    """
    Конвертирует PDF в список PNG-изображений (по одному на каждую страницу).

    Возвращает [(png_bytes, "image/png"), ...].
    Бросает ProcessingError при невалидном файле или превышении лимита страниц.
    """
    try:
        import fitz  # pymupdf
    except ImportError:
        raise ProcessingError("Сервер не поддерживает конвертацию PDF (pymupdf не установлен)")

    try:
        doc = fitz.open(stream=content, filetype="pdf")
    except Exception:
        raise ProcessingError("Не удалось открыть PDF. Проверьте, что файл не повреждён и не зашифрован")

    page_count = len(doc)
    if page_count == 0:
        doc.close()
        raise ProcessingError("PDF не содержит страниц")

    if page_count > max_pages:
        doc.close()
        raise ProcessingError(f"PDF содержит {page_count} страниц, максимум {max_pages}")

    pages = []
    try:
        for page in doc:
            pixmap = page.get_pixmap(dpi=150)
            pages.append((pixmap.tobytes("png"), "image/png"))
    except Exception as e:
        raise ProcessingError(f"Ошибка при рендеринге страницы PDF: {e}")
    finally:
        doc.close()

    return pages


def docx_to_text(content: bytes) -> tuple[bytes, str]:
    """
    Извлекает текст из DOCX и возвращает (text_bytes, "text/plain").
    Бросает ProcessingError при невалидном файле.
    """
    try:
        from docx import Document
    except ImportError:
        raise ProcessingError("Сервер не поддерживает конвертацию DOCX (python-docx не установлен)")

    try:
        doc = Document(io.BytesIO(content))
    except Exception:
        raise ProcessingError("Не удалось открыть DOCX. Проверьте, что файл не повреждён")

    lines = [para.text for para in doc.paragraphs]
    text = "\n".join(lines)

    return text.encode("utf-8"), "text/plain"


def pptx_to_text(content: bytes) -> tuple[bytes, str]:
    """
    Извлекает текст из всех слайдов PPTX и возвращает (text_bytes, "text/plain").
    Бросает ProcessingError при невалидном файле.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ProcessingError("Сервер не поддерживает конвертацию PPTX (python-pptx не установлен)")

    try:
        prs = Presentation(io.BytesIO(content))
    except Exception:
        raise ProcessingError("Не удалось открыть PPTX. Проверьте, что файл не повреждён")

    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            lines.append(f"[Слайд {i}]")
            lines.extend(slide_texts)

    text = "\n".join(lines)
    return text.encode("utf-8"), "text/plain"


def xlsx_to_text(content: bytes) -> tuple[bytes, str]:
    """
    Извлекает текст из всех листов XLSX и возвращает (text_bytes, "text/plain").
    Бросает ProcessingError при невалидном файле.
    """
    try:
        import openpyxl
    except ImportError:
        raise ProcessingError("Сервер не поддерживает конвертацию XLSX (openpyxl не установлен)")

    try:
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    except Exception:
        raise ProcessingError("Не удалось открыть XLSX. Проверьте, что файл не повреждён")

    lines = []
    for sheet in wb.worksheets:
        lines.append(f"[Лист: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))

    wb.close()
    text = "\n".join(lines)
    return text.encode("utf-8"), "text/plain"


def xls_to_text(content: bytes) -> tuple[bytes, str]:
    """
    Извлекает текст из всех листов XLS и возвращает (text_bytes, "text/plain").
    Бросает ProcessingError при невалидном файле.
    """
    try:
        import xlrd
    except ImportError:
        raise ProcessingError("Сервер не поддерживает конвертацию XLS (xlrd не установлен)")

    try:
        wb = xlrd.open_workbook(file_contents=content)
    except Exception:
        raise ProcessingError("Не удалось открыть XLS. Проверьте, что файл не повреждён")

    lines = []
    for sheet in wb.sheets():
        lines.append(f"[Лист: {sheet.name}]")
        for row_idx in range(sheet.nrows):
            cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)
                     if sheet.cell_value(row_idx, col) != ""]
            if cells:
                lines.append("\t".join(cells))

    text = "\n".join(lines)
    return text.encode("utf-8"), "text/plain"
